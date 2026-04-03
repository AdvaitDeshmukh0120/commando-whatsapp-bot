import os
import re
import time
import requests
from bs4 import BeautifulSoup
from urllib.parse import urljoin, urlparse
from PyPDF2 import PdfReader
from pptx import Presentation

BASE_URL = "https://www.commandonetworks.com"
OUTPUT_DIR = "data"
SCRAPED_FILE = os.path.join(OUTPUT_DIR, "scraped_website_data.txt")
DOWNLOAD_DIR = os.path.join(OUTPUT_DIR, "downloads")

HEADERS = {
    "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
}

SEED_URLS = [
    "/", "/company/about_us", "/switches", "/switches/C3500", "/switches/C3000",
    "/switches/E3000", "/switches/E1300", "/switches/IE3000", "/switches/C2000",
    "/switches/E2000", "/switches/E1100", "/switches/E1000", "/switches/IE1000",
    "/switches/E300", "/switches/E200", "/switches/E100", "/wireless", "/routers",
    "/gateways", "/accessories", "/warranty", "/support", "/partners/partner_program",
    "/contact", "/news", "/downloads", "/catalog", "/portfolio",
]


def clean_text(text):
    text = re.sub(r'\n\s*\n', '\n\n', text)
    text = re.sub(r'[ \t]+', ' ', text)
    lines = [line.strip() for line in text.split('\n')]
    return '\n'.join(line for line in lines if line)


def scrape_page(url):
    try:
        print(f"  Scraping: {url}")
        resp = requests.get(url, headers=HEADERS, timeout=15)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, 'html.parser')
        for tag in soup.find_all(['script', 'style', 'nav', 'footer', 'noscript']):
            tag.decompose()
        main = soup.find('main') or soup.find('div', class_='content') or soup.body
        if not main:
            return ""
        return clean_text(main.get_text(separator='\n'))
    except Exception as e:
        print(f"  Error scraping {url}: {e}")
        return ""


def find_downloadable_links(url):
    links = []
    try:
        resp = requests.get(url, headers=HEADERS, timeout=15)
        soup = BeautifulSoup(resp.text, 'html.parser')
        for a in soup.find_all('a', href=True):
            href = a['href'].lower()
            if href.endswith('.pdf') or href.endswith('.pptx') or href.endswith('.ppt'):
                links.append(urljoin(url, a['href']))
    except Exception as e:
        print(f"  Error: {e}")
    return list(set(links))


def download_file(url, dest_dir):
    try:
        filename = os.path.basename(urlparse(url).path)
        if not filename:
            return None
        filepath = os.path.join(dest_dir, filename)
        if os.path.exists(filepath):
            return filepath
        print(f"  Downloading: {filename}")
        resp = requests.get(url, headers=HEADERS, timeout=30)
        resp.raise_for_status()
        with open(filepath, 'wb') as f:
            f.write(resp.content)
        return filepath
    except Exception as e:
        print(f"  Error downloading {url}: {e}")
        return None


def extract_pdf_text(filepath):
    try:
        reader = PdfReader(filepath)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return clean_text(text)
    except Exception as e:
        print(f"  Error reading PDF: {e}")
        return ""


def extract_pptx_text(filepath):
    try:
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text += paragraph.text + "\n"
        return clean_text(text)
    except Exception as e:
        print(f"  Error reading PPTX: {e}")
        return ""


def run_scraper():
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    os.makedirs(DOWNLOAD_DIR, exist_ok=True)

    all_content = []
    all_download_links = []

    print("=" * 60)
    print("COMMANDO Networks Web Scraper")
    print("=" * 60)

    print("\n[Step 1] Scraping web pages...")
    for path in SEED_URLS:
        url = BASE_URL + path
        content = scrape_page(url)
        if content:
            all_content.append(f"===== PAGE: {path} =====\nURL: {url}\n\n{content}\n")
        links = find_downloadable_links(url)
        all_download_links.extend(links)
        time.sleep(0.5)

    all_download_links = list(set(all_download_links))
    if all_download_links:
        print(f"\n[Step 2] Downloading {len(all_download_links)} files...")
        for link in all_download_links:
            filepath = download_file(link, DOWNLOAD_DIR)
            if filepath:
                if filepath.lower().endswith('.pdf'):
                    text = extract_pdf_text(filepath)
                    if text:
                        all_content.append(f"===== PDF: {os.path.basename(filepath)} =====\nSource: {link}\n\n{text}\n")
                elif filepath.lower().endswith(('.pptx', '.ppt')):
                    text = extract_pptx_text(filepath)
                    if text:
                        all_content.append(f"===== PPTX: {os.path.basename(filepath)} =====\nSource: {link}\n\n{text}\n")
            time.sleep(0.5)

    print(f"\n[Step 3] Saving to {SCRAPED_FILE}...")
    with open(SCRAPED_FILE, 'w', encoding='utf-8') as f:
        f.write('\n\n'.join(all_content))

    print(f"\nDone! Scraped {len(SEED_URLS)} pages. Output: {SCRAPED_FILE}")


if __name__ == "__main__":
    run_scraper()
